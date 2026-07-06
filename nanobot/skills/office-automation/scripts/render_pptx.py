"""Render a slide DSL into a pptx file."""

from __future__ import annotations

import argparse
from pathlib import Path
from typing import Any

from _common import load_facts, read_json, render_text_value, replace_fact_placeholders
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def _add_textbox(slide: Any, left: float, top: float, width: float, height: float, text: str, size: int) -> Any:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(size)
    return box


def _add_title(slide: Any, title: str, subtitle: str | None = None) -> None:
    title_box = _add_textbox(slide, 0.45, 0.25, 9.1, 0.55, title, 28)
    title_box.text_frame.paragraphs[0].font.bold = True
    if subtitle:
        _add_textbox(slide, 0.5, 0.85, 9.0, 0.35, subtitle, 13)


def _add_bullets(slide: Any, bullets: list[Any], facts: dict[str, dict[str, Any]], top: float = 1.35) -> None:
    box = slide.shapes.add_textbox(Inches(0.75), Inches(top), Inches(8.7), Inches(4.7))
    frame = box.text_frame
    frame.clear()
    for index, item in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        text = item.get("text") if isinstance(item, dict) else item
        paragraph.text = replace_fact_placeholders(str(text), facts)
        paragraph.level = 0
        paragraph.font.size = Pt(18)


def _add_metrics(slide: Any, metrics: list[Any], facts: dict[str, dict[str, Any]]) -> None:
    if not metrics:
        return
    box_width = 2.8
    gap = 0.25
    start_left = 0.65
    top = 1.55
    for index, metric in enumerate(metrics[:3]):
        if not isinstance(metric, dict):
            continue
        left = start_left + index * (box_width + gap)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(box_width),
            Inches(1.25),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(232, 240, 254)
        shape.line.color.rgb = RGBColor(72, 96, 130)
        frame = shape.text_frame
        frame.clear()
        label = frame.paragraphs[0]
        label.text = str(metric.get("label", metric.get("fact_ref", "")))
        label.font.size = Pt(12)
        label.alignment = PP_ALIGN.CENTER
        value = frame.add_paragraph()
        value.text = render_text_value({"fact_ref": metric.get("fact_ref")}, facts)
        value.font.size = Pt(22)
        value.font.bold = True
        value.alignment = PP_ALIGN.CENTER


def _write_notes(slide: Any, text: str, facts: dict[str, dict[str, Any]]) -> None:
    if not text:
        return
    notes = slide.notes_slide.notes_text_frame
    notes.text = replace_fact_placeholders(text, facts)


def render_pptx(
    *,
    dsl_path: Path,
    facts_path: Path,
    output_path: Path,
    template_path: Path | None = None,
    constraints_path: Path | None = None,
) -> None:
    dsl = read_json(dsl_path)
    facts = load_facts(facts_path)
    constraints = read_json(constraints_path) if constraints_path else None
    slides_payload = dsl.get("slides", [])
    if not isinstance(slides_payload, list) or not slides_payload:
        raise ValueError("slide DSL must contain slides")

    max_pages = constraints.get("outputs", {}).get("pptx_max_pages") if constraints else None
    if isinstance(max_pages, int) and len(slides_payload) > max_pages:
        raise ValueError(f"slide count {len(slides_payload)} exceeds limit {max_pages}")

    presentation = Presentation(str(template_path)) if template_path else Presentation()
    blank_layout = presentation.slide_layouts[6]
    while len(presentation.slides) > 0:
        xml_slides = presentation.slides._sldIdLst  # noqa: SLF001 - python-pptx has no public clear API.
        rel_id = xml_slides[0].rId
        presentation.part.drop_rel(rel_id)
        xml_slides.remove(xml_slides[0])

    for slide_payload in slides_payload:
        if not isinstance(slide_payload, dict):
            continue
        title = replace_fact_placeholders(str(slide_payload.get("title", "")), facts)
        if not title:
            raise ValueError("each slide must have a title")
        slide = presentation.slides.add_slide(blank_layout)
        subtitle = slide_payload.get("subtitle")
        _add_title(slide, title, replace_fact_placeholders(str(subtitle), facts) if subtitle else None)
        metrics = slide_payload.get("metrics", [])
        if isinstance(metrics, list):
            _add_metrics(slide, metrics, facts)
        bullets = slide_payload.get("bullets", [])
        if isinstance(bullets, list):
            _add_bullets(slide, bullets, facts, top=3.15 if metrics else 1.35)
        notes = slide_payload.get("speaker_notes")
        if isinstance(notes, str):
            _write_notes(slide, notes, facts)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--dsl", dest="dsl_path", required=True, type=Path)
    parser.add_argument("--facts", dest="facts_path", required=True, type=Path)
    parser.add_argument("--template", dest="template_path", type=Path)
    parser.add_argument("--constraints", dest="constraints_path", type=Path)
    parser.add_argument("--out", dest="output_path", required=True, type=Path)
    args = parser.parse_args()

    render_pptx(
        dsl_path=args.dsl_path,
        facts_path=args.facts_path,
        output_path=args.output_path,
        template_path=args.template_path,
        constraints_path=args.constraints_path,
    )


if __name__ == "__main__":
    main()
