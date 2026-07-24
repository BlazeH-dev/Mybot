"""OfficePython request/result JSON interface for DOCX, XLSX, and PPTX artifacts."""

from __future__ import annotations

import argparse
import hashlib
import json
import os
import subprocess
import sys
import tempfile
import time
import zipfile
from collections.abc import Callable
from datetime import date, datetime
from pathlib import Path
from typing import Any

from docx import Document
from docx.shared import Inches as DocxInches
from openpyxl import Workbook, load_workbook
from openpyxl.chart import BarChart, LineChart, Reference
from pptx import Presentation
from pptx.util import Inches as PptxInches

SCHEMA_VERSION = 1
FORMATS = {"docx", "xlsx", "pptx"}
OPERATIONS = {"inspect", "query", "create", "apply", "validate", "render"}


class OfficeRequestError(Exception):
    """Structured request failure returned through the result JSON file."""

    def __init__(self, code: str, message: str, *, status: str = "error") -> None:
        super().__init__(message)
        self.code = code
        self.status = status


def _json_value(value: Any) -> Any:
    if isinstance(value, (date, datetime)):
        return value.isoformat()
    if isinstance(value, (str, int, float, bool)) or value is None:
        return value
    return str(value)


def _read_json(path: Path) -> dict[str, Any]:
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        raise OfficeRequestError("invalid_request_json", str(exc)) from exc
    if not isinstance(payload, dict):
        raise OfficeRequestError("invalid_request", "request must be a JSON object")
    return payload


def _write_json_atomic(path: Path, payload: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    handle, temporary = tempfile.mkstemp(prefix=f".{path.name}.", suffix=".tmp", dir=path.parent)
    try:
        with os.fdopen(handle, "w", encoding="utf-8") as stream:
            json.dump(payload, stream, ensure_ascii=False, indent=2)
            stream.write("\n")
            stream.flush()
            os.fsync(stream.fileno())
        os.replace(temporary, path)
    except Exception:
        Path(temporary).unlink(missing_ok=True)
        raise


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _artifact(path: Path) -> dict[str, Any]:
    return {
        "path": str(path),
        "sha256": _sha256(path),
        "size_bytes": path.stat().st_size,
    }


def _artifact_path(value: Any, *, field: str, required: bool = True) -> Path | None:
    if value is None and not required:
        return None
    if not isinstance(value, dict) or not isinstance(value.get("path"), str):
        raise OfficeRequestError("invalid_artifact", f"{field}.path is required")
    return Path(value["path"]).expanduser().resolve()


def _within(path: Path, root: Path) -> bool:
    try:
        path.relative_to(root)
    except ValueError:
        return False
    return True


def _validate_request(request: dict[str, Any]) -> tuple[str, str, Path, Path | None, Path]:
    if request.get("schema_version") != SCHEMA_VERSION:
        raise OfficeRequestError(
            "unsupported_schema",
            f"schema_version must be {SCHEMA_VERSION}",
            status="unsupported",
        )
    operation = request.get("operation")
    office_format = request.get("format")
    if operation not in OPERATIONS:
        raise OfficeRequestError("unsupported_operation", f"unsupported operation: {operation}", status="unsupported")
    if office_format not in FORMATS:
        raise OfficeRequestError("unsupported_format", f"unsupported format: {office_format}", status="unsupported")

    options = request.get("options")
    if not isinstance(options, dict):
        raise OfficeRequestError("invalid_options", "options must be an object")
    root_value = options.get("artifact_root")
    if not isinstance(root_value, str) or not root_value:
        raise OfficeRequestError("missing_artifact_root", "options.artifact_root is required")
    artifact_root = Path(root_value).expanduser().resolve()

    input_path = _artifact_path(
        request.get("input_artifact"),
        field="input_artifact",
        required=operation != "create",
    )
    output_path = _artifact_path(
        request.get("output_artifact"),
        field="output_artifact",
        required=operation in {"create", "apply", "render"},
    )
    if output_path is not None:
        if not _within(output_path, artifact_root):
            raise OfficeRequestError(
                "output_outside_artifact_root",
                "output_artifact must be inside options.artifact_root",
            )
        if input_path is not None and output_path == input_path:
            raise OfficeRequestError("readonly_input", "input and output artifacts must differ")
        expected_suffix = ".pdf" if operation == "render" else f".{office_format}"
        if output_path.suffix.lower() != expected_suffix:
            raise OfficeRequestError(
                "invalid_output_extension",
                f"{operation} output must use {expected_suffix}",
            )
    if input_path is not None:
        if not input_path.is_file():
            raise OfficeRequestError("input_not_found", f"input artifact not found: {input_path}")
        if input_path.suffix.lower() != f".{office_format}":
            raise OfficeRequestError(
                "format_mismatch",
                f"input extension does not match format {office_format}",
            )
    return operation, office_format, artifact_root, input_path, output_path or artifact_root


def _base_result(operation: Any = None, office_format: Any = None) -> dict[str, Any]:
    return {
        "schema_version": SCHEMA_VERSION,
        "status": "ok",
        "operation": operation,
        "format": office_format,
        "matches": [],
        "changes": [],
        "artifact": None,
        "validation": None,
        "rendered_assets": [],
        "warnings": [],
        "error": None,
    }


def _atomic_office_write(
    output_path: Path,
    office_format: str,
    writer: Callable[[Path], None],
) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    handle, temporary = tempfile.mkstemp(
        prefix=f".{output_path.stem}.",
        suffix=output_path.suffix,
        dir=output_path.parent,
    )
    os.close(handle)
    temp_path = Path(temporary)
    try:
        writer(temp_path)
        if not temp_path.is_file() or temp_path.stat().st_size == 0:
            raise OfficeRequestError("empty_output", "Office writer produced no output")
        validation = _validate_openxml(temp_path, office_format)
        if not validation["passed"]:
            messages = "; ".join(issue["message"] for issue in validation["issues"])
            raise OfficeRequestError("validation_failed", messages or "Office validation failed")
        os.replace(temp_path, output_path)
    finally:
        temp_path.unlink(missing_ok=True)


def _zip_member_text(path: Path, member: str) -> str:
    try:
        with zipfile.ZipFile(path) as archive:
            return archive.read(member).decode("utf-8", errors="replace")
    except (KeyError, OSError, zipfile.BadZipFile):
        return ""


def _unsupported_features(path: Path, office_format: str) -> list[str]:
    features: list[str] = []
    try:
        with zipfile.ZipFile(path) as archive:
            names = set(archive.namelist())
            if office_format == "docx":
                document_xml = _zip_member_text(path, "word/document.xml")
                settings_xml = _zip_member_text(path, "word/settings.xml")
                if any(token in document_xml for token in ("<w:ins", "<w:del", "<w:moveFrom", "<w:moveTo")):
                    features.append("tracked_changes")
                if "<w:trackRevisions" in settings_xml:
                    features.append("tracked_changes_enabled")
            elif office_format == "pptx":
                if any(name.startswith("ppt/diagrams/") for name in names):
                    features.append("smartart")
                if len([name for name in names if name.startswith("ppt/slideMasters/slideMaster") and name.endswith(".xml")]) > 1:
                    features.append("complex_master_set")
                for name in names:
                    if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                        if "<p:timing" in _zip_member_text(path, name):
                            features.append("animations_or_timing")
                            break
    except (OSError, zipfile.BadZipFile) as exc:
        raise OfficeRequestError("invalid_openxml", str(exc)) from exc
    return sorted(set(features))


def _validate_openxml(path: Path, office_format: str) -> dict[str, Any]:
    issues: list[dict[str, str]] = []
    try:
        if office_format == "docx":
            document = Document(path)
            if not document.paragraphs and not document.tables:
                issues.append({"code": "empty_document", "message": "document has no body content"})
        elif office_format == "xlsx":
            workbook = load_workbook(path, data_only=False, read_only=True)
            if not workbook.sheetnames:
                issues.append({"code": "empty_workbook", "message": "workbook has no worksheets"})
            workbook.close()
        else:
            presentation = Presentation(path)
            if not presentation.slides:
                issues.append({"code": "empty_presentation", "message": "presentation has no slides"})
    except Exception as exc:
        issues.append({"code": "open_failed", "message": str(exc)})
    return {"passed": not issues, "issues": issues, "unsupported_features": _unsupported_features(path, office_format)}


def _paragraph_payload(paragraph: Any, index: int, *, scope: str = "body") -> dict[str, Any]:
    style = paragraph.style.name if paragraph.style is not None else None
    return {"kind": "paragraph", "scope": scope, "index": index, "text": paragraph.text, "style": style}


def _docx_matches(document: Document, selector: dict[str, Any]) -> list[dict[str, Any]]:
    kind = selector.get("kind", "document")
    matches: list[dict[str, Any]] = []
    if kind in {"document", "paragraph"}:
        for index, paragraph in enumerate(document.paragraphs):
            matches.append(_paragraph_payload(paragraph, index))
    if kind in {"document", "table", "cell"}:
        for table_index, table in enumerate(document.tables):
            if kind in {"document", "table"}:
                matches.append(
                    {
                        "kind": "table",
                        "index": table_index,
                        "rows": len(table.rows),
                        "columns": len(table.columns),
                        "values": [[cell.text for cell in row.cells] for row in table.rows],
                    }
                )
            if kind == "cell":
                for row_index, row in enumerate(table.rows):
                    for column_index, cell in enumerate(row.cells):
                        matches.append(
                            {
                                "kind": "cell",
                                "table_index": table_index,
                                "row": row_index,
                                "column": column_index,
                                "text": cell.text,
                            }
                        )
    if kind in {"header", "footer"}:
        for section_index, section in enumerate(document.sections):
            container = section.header if kind == "header" else section.footer
            for index, paragraph in enumerate(container.paragraphs):
                item = _paragraph_payload(paragraph, index, scope=kind)
                item["section"] = section_index
                matches.append(item)
    if kind == "image":
        matches = [
            {"kind": "image", "index": index, "relationship_id": shape._inline.graphic.graphicData.pic.blipFill.blip.embed}
            for index, shape in enumerate(document.inline_shapes)
        ]
    return _filter_matches(matches, selector)


def _filter_matches(matches: list[dict[str, Any]], selector: dict[str, Any]) -> list[dict[str, Any]]:
    index = selector.get("index")
    name = selector.get("name")
    text_contains = selector.get("text_contains")
    filtered: list[dict[str, Any]] = []
    for match in matches:
        if index is not None and match.get("index") != index:
            continue
        if name is not None and match.get("name") != name:
            continue
        if text_contains is not None and str(text_contains) not in str(match.get("text", "")):
            continue
        filtered.append(match)
    return filtered


def _inspect_docx(path: Path) -> dict[str, Any]:
    document = Document(path)
    return {
        "paragraph_count": len(document.paragraphs),
        "table_count": len(document.tables),
        "image_count": len(document.inline_shapes),
        "section_count": len(document.sections),
        "styles": sorted({p.style.name for p in document.paragraphs if p.style is not None}),
        "unsupported_features": _unsupported_features(path, "docx"),
    }


def _create_docx(payload: dict[str, Any], output_path: Path) -> list[dict[str, Any]]:
    def writer(temp_path: Path) -> None:
        document = Document()
        title = payload.get("title")
        if isinstance(title, str) and title:
            document.add_heading(title, level=0)
        for item in payload.get("paragraphs", []):
            if isinstance(item, str):
                document.add_paragraph(item)
            elif isinstance(item, dict):
                text = str(item.get("text", ""))
                style = item.get("style")
                if item.get("heading_level") is not None:
                    document.add_heading(text, level=int(item["heading_level"]))
                else:
                    document.add_paragraph(text, style=str(style) if style else None)
        for item in payload.get("lists", []):
            if not isinstance(item, dict):
                continue
            style = "List Number" if item.get("ordered") else "List Bullet"
            for value in item.get("items", []):
                document.add_paragraph(str(value), style=style)
        for table_payload in payload.get("tables", []):
            rows = table_payload.get("rows", []) if isinstance(table_payload, dict) else []
            if not rows:
                continue
            width = max(len(row) for row in rows if isinstance(row, list))
            table = document.add_table(rows=0, cols=width)
            table.style = str(table_payload.get("style", "Table Grid"))
            for values in rows:
                cells = table.add_row().cells
                for index, value in enumerate(values[:width]):
                    cells[index].text = str(value)
        for image in payload.get("images", []):
            if not isinstance(image, dict) or not isinstance(image.get("path"), str):
                raise OfficeRequestError("invalid_image", "DOCX image.path is required")
            document.add_picture(image["path"], width=DocxInches(float(image.get("width_inches", 4))))
        if isinstance(payload.get("header"), str):
            document.sections[0].header.paragraphs[0].text = payload["header"]
        if isinstance(payload.get("footer"), str):
            document.sections[0].footer.paragraphs[0].text = payload["footer"]
        document.save(temp_path)

    _atomic_office_write(output_path, "docx", writer)
    return [{"type": "create", "path": str(output_path)}]


def _selected_docx_objects(document: Document, selector: dict[str, Any]) -> list[Any]:
    kind = selector.get("kind", "paragraph")
    matches = _docx_matches(document, selector)
    selected: list[Any] = []
    for match in matches:
        if kind == "paragraph":
            selected.append(document.paragraphs[match["index"]])
        elif kind == "cell":
            selected.append(
                document.tables[match["table_index"]].rows[match["row"]].cells[match["column"]]
            )
        elif kind in {"header", "footer"}:
            section = document.sections[match["section"]]
            container = section.header if kind == "header" else section.footer
            selected.append(container.paragraphs[match["index"]])
    return selected


def _apply_docx(source: Path, payload: dict[str, Any], output_path: Path) -> list[dict[str, Any]]:
    unsupported = _unsupported_features(source, "docx")
    if unsupported:
        raise OfficeRequestError(
            "unsupported_features",
            f"cannot safely edit DOCX containing: {', '.join(unsupported)}",
            status="unsupported",
        )
    actions = payload.get("actions")
    if not isinstance(actions, list) or not actions:
        raise OfficeRequestError("missing_actions", "payload.actions must be a non-empty array")
    changes: list[dict[str, Any]] = []

    def writer(temp_path: Path) -> None:
        document = Document(source)
        for action in actions:
            if not isinstance(action, dict):
                raise OfficeRequestError("invalid_action", "each action must be an object")
            action_type = action.get("type")
            selector = action.get("selector", {})
            if action_type in {"set_text", "set_style"}:
                selected = _selected_docx_objects(document, selector)
                if not selected:
                    raise OfficeRequestError("selector_no_match", f"selector matched no DOCX objects: {selector}")
                for item in selected:
                    if action_type == "set_text":
                        item.text = str(action.get("value", ""))
                    else:
                        item.style = str(action.get("value", "Normal"))
                changes.append({"type": action_type, "count": len(selected), "selector": selector})
            elif action_type == "append_paragraph":
                document.add_paragraph(str(action.get("text", "")), style=action.get("style"))
                changes.append({"type": action_type, "count": 1})
            elif action_type == "append_list":
                items = action.get("items", [])
                style = "List Number" if action.get("ordered") else "List Bullet"
                for item in items:
                    document.add_paragraph(str(item), style=style)
                changes.append({"type": action_type, "count": len(items)})
            elif action_type == "append_table":
                rows = action.get("rows", [])
                if not rows:
                    raise OfficeRequestError("invalid_action", "append_table.rows is required")
                width = max(len(row) for row in rows)
                table = document.add_table(rows=0, cols=width)
                table.style = str(action.get("style", "Table Grid"))
                for values in rows:
                    cells = table.add_row().cells
                    for index, value in enumerate(values[:width]):
                        cells[index].text = str(value)
                changes.append({"type": action_type, "count": len(rows)})
            elif action_type == "add_image":
                document.add_picture(
                    str(action.get("path", "")),
                    width=DocxInches(float(action.get("width_inches", 4))),
                )
                changes.append({"type": action_type, "count": 1})
            elif action_type in {"set_header", "set_footer"}:
                for section in document.sections:
                    target = section.header if action_type == "set_header" else section.footer
                    target.paragraphs[0].text = str(action.get("text", ""))
                changes.append({"type": action_type, "count": len(document.sections)})
            else:
                raise OfficeRequestError("unsupported_action", f"unsupported DOCX action: {action_type}", status="unsupported")
        document.save(temp_path)

    _atomic_office_write(output_path, "docx", writer)
    return changes


def _xlsx_sheet(workbook: Any, selector: dict[str, Any]) -> Any:
    name = selector.get("sheet") or selector.get("name")
    if name is None:
        return workbook.active
    if name not in workbook.sheetnames:
        raise OfficeRequestError("selector_no_match", f"worksheet not found: {name}")
    return workbook[name]


def _xlsx_matches(workbook: Any, selector: dict[str, Any]) -> list[dict[str, Any]]:
    kind = selector.get("kind", "workbook")
    if kind in {"workbook", "sheet"}:
        matches = [
            {
                "kind": "sheet",
                "index": index,
                "name": sheet.title,
                "rows": sheet.max_row,
                "columns": sheet.max_column,
            }
            for index, sheet in enumerate(workbook.worksheets)
        ]
        return _filter_matches(matches, selector)
    if kind in {"cell", "formula"}:
        sheet = _xlsx_sheet(workbook, selector)
        cell_range = selector.get("range")
        cells = sheet[cell_range] if isinstance(cell_range, str) else sheet.iter_rows()
        if hasattr(cells, "coordinate"):
            cells = ((cells,),)
        elif isinstance(cells, tuple) and cells and not isinstance(cells[0], tuple):
            cells = (cells,)
        matches = []
        for row in cells:
            for cell in row:
                if kind == "formula" and not (isinstance(cell.value, str) and cell.value.startswith("=")):
                    continue
                matches.append(
                    {
                        "kind": "cell",
                        "sheet": sheet.title,
                        "coordinate": cell.coordinate,
                        "value": _json_value(cell.value),
                        "data_type": cell.data_type,
                    }
                )
        return matches
    if kind == "chart":
        sheet = _xlsx_sheet(workbook, selector)
        return [
            {"kind": "chart", "index": index, "title": str(chart.title), "sheet": sheet.title}
            for index, chart in enumerate(sheet._charts)
        ]
    raise OfficeRequestError("unsupported_selector", f"unsupported XLSX selector kind: {kind}", status="unsupported")


def _inspect_xlsx(path: Path) -> dict[str, Any]:
    workbook = load_workbook(path, data_only=False, read_only=False)
    try:
        return {
            "sheets": _xlsx_matches(workbook, {"kind": "sheet"}),
            "formula_count": sum(
                1
                for sheet in workbook.worksheets
                for row in sheet.iter_rows()
                for cell in row
                if isinstance(cell.value, str) and cell.value.startswith("=")
            ),
            "chart_count": sum(len(sheet._charts) for sheet in workbook.worksheets),
            "defined_names": sorted(item.name for item in workbook.defined_names.values()),
            "unsupported_features": [],
        }
    finally:
        workbook.close()


def _add_xlsx_chart(sheet: Any, spec: dict[str, Any]) -> None:
    chart_type = spec.get("chart_type", "bar")
    if chart_type == "bar":
        chart = BarChart()
    elif chart_type == "line":
        chart = LineChart()
    else:
        raise OfficeRequestError("unsupported_chart", f"unsupported chart type: {chart_type}", status="unsupported")
    data = spec.get("data")
    if not isinstance(data, dict):
        raise OfficeRequestError("invalid_chart", "chart data range is required")
    chart.add_data(
        Reference(
            sheet,
            min_col=int(data["min_col"]),
            max_col=int(data.get("max_col", data["min_col"])),
            min_row=int(data["min_row"]),
            max_row=int(data["max_row"]),
        ),
        titles_from_data=bool(data.get("titles_from_data", True)),
    )
    categories = spec.get("categories")
    if isinstance(categories, dict):
        chart.set_categories(
            Reference(
                sheet,
                min_col=int(categories["min_col"]),
                min_row=int(categories["min_row"]),
                max_row=int(categories["max_row"]),
            )
        )
    chart.title = str(spec.get("title", ""))
    sheet.add_chart(chart, str(spec.get("anchor", "E2")))


def _populate_xlsx(workbook: Any, payload: dict[str, Any]) -> None:
    sheets = payload.get("sheets", [])
    if not isinstance(sheets, list) or not sheets:
        raise OfficeRequestError("missing_sheets", "payload.sheets must be a non-empty array")
    while workbook.worksheets:
        workbook.remove(workbook.worksheets[0])
    for sheet_payload in sheets:
        if not isinstance(sheet_payload, dict) or not isinstance(sheet_payload.get("name"), str):
            raise OfficeRequestError("invalid_sheet", "each sheet requires a name")
        sheet = workbook.create_sheet(sheet_payload["name"])
        for row in sheet_payload.get("rows", []):
            sheet.append([_json_value(value) for value in row])
        for coordinate, value in sheet_payload.get("cells", {}).items():
            sheet[str(coordinate)] = value
        for coordinate, formula in sheet_payload.get("formulas", {}).items():
            sheet[str(coordinate)] = str(formula)
        for chart in sheet_payload.get("charts", []):
            _add_xlsx_chart(sheet, chart)


def _create_xlsx(payload: dict[str, Any], output_path: Path) -> list[dict[str, Any]]:
    def writer(temp_path: Path) -> None:
        workbook = Workbook()
        _populate_xlsx(workbook, payload)
        workbook.save(temp_path)
        workbook.close()

    _atomic_office_write(output_path, "xlsx", writer)
    return [{"type": "create", "path": str(output_path)}]


def _apply_xlsx(source: Path, payload: dict[str, Any], output_path: Path) -> list[dict[str, Any]]:
    actions = payload.get("actions")
    if not isinstance(actions, list) or not actions:
        raise OfficeRequestError("missing_actions", "payload.actions must be a non-empty array")
    changes: list[dict[str, Any]] = []

    def writer(temp_path: Path) -> None:
        workbook = load_workbook(source, data_only=False)
        try:
            for action in actions:
                if not isinstance(action, dict):
                    raise OfficeRequestError("invalid_action", "each action must be an object")
                action_type = action.get("type")
                selector = action.get("selector", {})
                if action_type in {"set_cell", "clear_cell"}:
                    sheet = _xlsx_sheet(workbook, selector)
                    coordinate = selector.get("coordinate")
                    if not isinstance(coordinate, str):
                        raise OfficeRequestError("invalid_selector", "cell selector.coordinate is required")
                    sheet[coordinate] = None if action_type == "clear_cell" else action.get("value")
                    changes.append({"type": action_type, "count": 1, "sheet": sheet.title, "coordinate": coordinate})
                elif action_type == "append_row":
                    sheet = _xlsx_sheet(workbook, selector)
                    values = action.get("values", [])
                    sheet.append(values)
                    changes.append({"type": action_type, "count": 1, "sheet": sheet.title})
                elif action_type == "add_sheet":
                    name = action.get("name")
                    if not isinstance(name, str) or name in workbook.sheetnames:
                        raise OfficeRequestError("invalid_sheet", f"invalid or duplicate sheet name: {name}")
                    sheet = workbook.create_sheet(name)
                    for row in action.get("rows", []):
                        sheet.append(row)
                    changes.append({"type": action_type, "count": 1, "sheet": name})
                elif action_type == "rename_sheet":
                    sheet = _xlsx_sheet(workbook, selector)
                    name = action.get("name")
                    if not isinstance(name, str) or name in workbook.sheetnames:
                        raise OfficeRequestError("invalid_sheet", f"invalid or duplicate sheet name: {name}")
                    old_name = sheet.title
                    sheet.title = name
                    changes.append({"type": action_type, "count": 1, "from": old_name, "to": name})
                elif action_type == "add_chart":
                    sheet = _xlsx_sheet(workbook, selector)
                    _add_xlsx_chart(sheet, action)
                    changes.append({"type": action_type, "count": 1, "sheet": sheet.title})
                else:
                    raise OfficeRequestError("unsupported_action", f"unsupported XLSX action: {action_type}", status="unsupported")
            workbook.save(temp_path)
        finally:
            workbook.close()

    _atomic_office_write(output_path, "xlsx", writer)
    return changes


def _pptx_matches(presentation: Presentation, selector: dict[str, Any]) -> list[dict[str, Any]]:
    kind = selector.get("kind", "presentation")
    matches: list[dict[str, Any]] = []
    for slide_index, slide in enumerate(presentation.slides):
        if kind in {"presentation", "slide"}:
            title = slide.shapes.title.text if slide.shapes.title is not None else ""
            matches.append({"kind": "slide", "index": slide_index, "title": title, "shape_count": len(slide.shapes)})
        if kind in {"shape", "text"}:
            for shape_index, shape in enumerate(slide.shapes):
                text = shape.text if hasattr(shape, "text") else ""
                matches.append(
                    {
                        "kind": "shape",
                        "index": shape_index,
                        "slide": slide_index,
                        "name": shape.name,
                        "shape_type": str(shape.shape_type),
                        "text": text,
                    }
                )
    slide_filter = selector.get("slide")
    if slide_filter is not None:
        matches = [match for match in matches if match.get("slide", match.get("index")) == slide_filter]
    return _filter_matches(matches, selector)


def _inspect_pptx(path: Path) -> dict[str, Any]:
    presentation = Presentation(path)
    return {
        "slide_count": len(presentation.slides),
        "slides": _pptx_matches(presentation, {"kind": "slide"}),
        "layout_count": len(presentation.slide_layouts),
        "unsupported_features": _unsupported_features(path, "pptx"),
    }


def _add_pptx_slide(presentation: Presentation, spec: dict[str, Any]) -> None:
    layout_kind = spec.get("layout", "title_and_content")
    if layout_kind == "title":
        layout_index = 0
    elif layout_kind == "title_and_content":
        layout_index = 1
    elif layout_kind == "blank":
        layout_index = 6
    else:
        raise OfficeRequestError("unsupported_layout", f"unsupported PPTX layout: {layout_kind}", status="unsupported")
    slide = presentation.slides.add_slide(presentation.slide_layouts[layout_index])
    if slide.shapes.title is not None:
        slide.shapes.title.text = str(spec.get("title", ""))
    body = spec.get("body")
    if body is not None and len(slide.placeholders) > 1:
        frame = slide.placeholders[1].text_frame
        frame.clear()
        values = body if isinstance(body, list) else [body]
        for index, value in enumerate(values):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = str(value)
    for text_box in spec.get("text_boxes", []):
        box = slide.shapes.add_textbox(
            PptxInches(float(text_box.get("left", 1))),
            PptxInches(float(text_box.get("top", 1))),
            PptxInches(float(text_box.get("width", 4))),
            PptxInches(float(text_box.get("height", 1))),
        )
        box.text_frame.text = str(text_box.get("text", ""))
    for image in spec.get("images", []):
        slide.shapes.add_picture(
            str(image.get("path", "")),
            PptxInches(float(image.get("left", 1))),
            PptxInches(float(image.get("top", 1))),
            width=PptxInches(float(image.get("width", 4))),
        )


def _create_pptx(payload: dict[str, Any], output_path: Path) -> list[dict[str, Any]]:
    slides = payload.get("slides")
    if not isinstance(slides, list) or not slides:
        raise OfficeRequestError("missing_slides", "payload.slides must be a non-empty array")

    def writer(temp_path: Path) -> None:
        presentation = Presentation()
        for slide in slides:
            if not isinstance(slide, dict):
                raise OfficeRequestError("invalid_slide", "each slide must be an object")
            _add_pptx_slide(presentation, slide)
        presentation.save(temp_path)

    _atomic_office_write(output_path, "pptx", writer)
    return [{"type": "create", "path": str(output_path), "count": len(slides)}]


def _apply_pptx(source: Path, payload: dict[str, Any], output_path: Path) -> list[dict[str, Any]]:
    unsupported = _unsupported_features(source, "pptx")
    if unsupported:
        raise OfficeRequestError(
            "unsupported_features",
            f"cannot safely edit PPTX containing: {', '.join(unsupported)}",
            status="unsupported",
        )
    actions = payload.get("actions")
    if not isinstance(actions, list) or not actions:
        raise OfficeRequestError("missing_actions", "payload.actions must be a non-empty array")
    changes: list[dict[str, Any]] = []

    def writer(temp_path: Path) -> None:
        presentation = Presentation(source)
        for action in actions:
            if not isinstance(action, dict):
                raise OfficeRequestError("invalid_action", "each action must be an object")
            action_type = action.get("type")
            selector = action.get("selector", {})
            if action_type == "set_text":
                matches = _pptx_matches(presentation, selector)
                if not matches:
                    raise OfficeRequestError("selector_no_match", f"selector matched no PPTX shapes: {selector}")
                for match in matches:
                    if match["kind"] != "shape":
                        raise OfficeRequestError("invalid_selector", "set_text requires a shape selector")
                    shape = presentation.slides[match["slide"]].shapes[match["index"]]
                    if not hasattr(shape, "text_frame"):
                        raise OfficeRequestError("unsupported_shape", f"shape cannot contain text: {shape.name}", status="unsupported")
                    shape.text = str(action.get("value", ""))
                changes.append({"type": action_type, "count": len(matches), "selector": selector})
            elif action_type == "add_slide":
                spec = action.get("slide")
                if not isinstance(spec, dict):
                    raise OfficeRequestError("invalid_action", "add_slide.slide is required")
                _add_pptx_slide(presentation, spec)
                changes.append({"type": action_type, "count": 1})
            elif action_type == "add_text":
                slide_index = int(selector.get("slide", 0))
                slide = presentation.slides[slide_index]
                box = slide.shapes.add_textbox(
                    PptxInches(float(action.get("left", 1))),
                    PptxInches(float(action.get("top", 1))),
                    PptxInches(float(action.get("width", 4))),
                    PptxInches(float(action.get("height", 1))),
                )
                box.text_frame.text = str(action.get("text", ""))
                changes.append({"type": action_type, "count": 1, "slide": slide_index})
            elif action_type == "add_image":
                slide_index = int(selector.get("slide", 0))
                presentation.slides[slide_index].shapes.add_picture(
                    str(action.get("path", "")),
                    PptxInches(float(action.get("left", 1))),
                    PptxInches(float(action.get("top", 1))),
                    width=PptxInches(float(action.get("width", 4))),
                )
                changes.append({"type": action_type, "count": 1, "slide": slide_index})
            else:
                raise OfficeRequestError("unsupported_action", f"unsupported PPTX action: {action_type}", status="unsupported")
        presentation.save(temp_path)

    _atomic_office_write(output_path, "pptx", writer)
    return changes


def _run_render(request: dict[str, Any], input_path: Path, output_path: Path) -> dict[str, Any]:
    options = request["options"]
    libreoffice = options.get("libreoffice")
    if not isinstance(libreoffice, dict):
        raise OfficeRequestError("libreoffice_not_configured", "options.libreoffice is required for render")
    executable_value = libreoffice.get("path")
    expected_version = libreoffice.get("expected_version")
    if not isinstance(executable_value, str) or not isinstance(expected_version, str):
        raise OfficeRequestError(
            "libreoffice_not_locked",
            "libreoffice.path and libreoffice.expected_version are required",
        )
    executable = Path(executable_value).expanduser().resolve()
    if not executable.is_file() or not os.access(executable, os.X_OK):
        raise OfficeRequestError("libreoffice_not_found", f"LibreOffice executable is unavailable: {executable}")
    timeout = float(libreoffice.get("timeout_seconds", 120))
    try:
        version_run = subprocess.run(
            [str(executable), "--version"],
            capture_output=True,
            text=True,
            timeout=min(timeout, 30),
            check=False,
        )
    except (OSError, subprocess.SubprocessError) as exc:
        raise OfficeRequestError("libreoffice_probe_failed", str(exc)) from exc
    actual_version = version_run.stdout.strip() or version_run.stderr.strip()
    if version_run.returncode != 0:
        raise OfficeRequestError("libreoffice_probe_failed", actual_version)
    if actual_version != expected_version:
        raise OfficeRequestError(
            "libreoffice_version_mismatch",
            f"expected {expected_version!r}, got {actual_version!r}",
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="office-python-render-", dir=output_path.parent) as temp_dir:
        try:
            render_run = subprocess.run(
                [
                    str(executable),
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    temp_dir,
                    str(input_path),
                ],
                capture_output=True,
                text=True,
                timeout=timeout,
                check=False,
            )
        except (OSError, subprocess.SubprocessError) as exc:
            raise OfficeRequestError("render_failed", str(exc)) from exc
        rendered = Path(temp_dir) / f"{input_path.stem}.pdf"
        if render_run.returncode != 0 or not rendered.is_file():
            detail = render_run.stderr.strip() or render_run.stdout.strip() or "no PDF produced"
            raise OfficeRequestError("render_failed", detail)
        os.replace(rendered, output_path)
    return {
        "artifact": _artifact(output_path),
        "rendered_assets": [_artifact(output_path)],
        "validation": {
            "passed": True,
            "engine": "libreoffice",
            "path": str(executable),
            "version": actual_version,
        },
    }


def execute(request: dict[str, Any]) -> dict[str, Any]:
    operation, office_format, _artifact_root, input_path, output_path = _validate_request(request)
    result = _base_result(operation, office_format)
    selector = request.get("selector")
    payload = request.get("payload")
    if not isinstance(selector, dict):
        raise OfficeRequestError("invalid_selector", "selector must be an object")
    if not isinstance(payload, dict):
        raise OfficeRequestError("invalid_payload", "payload must be an object")

    input_hash = _sha256(input_path) if input_path is not None else None
    if operation == "render":
        assert input_path is not None
        render_result = _run_render(request, input_path, output_path)
        result.update(render_result)
    elif operation in {"inspect", "query"}:
        assert input_path is not None
        if office_format == "docx":
            matches = [_inspect_docx(input_path)] if operation == "inspect" else _docx_matches(Document(input_path), selector)
        elif office_format == "xlsx":
            workbook = load_workbook(input_path, data_only=bool(request["options"].get("data_only", False)))
            try:
                matches = [_inspect_xlsx(input_path)] if operation == "inspect" else _xlsx_matches(workbook, selector)
            finally:
                workbook.close()
        else:
            presentation = Presentation(input_path)
            matches = [_inspect_pptx(input_path)] if operation == "inspect" else _pptx_matches(presentation, selector)
        result["matches"] = matches
    elif operation == "validate":
        assert input_path is not None
        validation = _validate_openxml(input_path, office_format)
        result["validation"] = validation
        if validation["unsupported_features"]:
            result["status"] = "unsupported"
            result["error"] = {
                "code": "unsupported_features",
                "message": ", ".join(validation["unsupported_features"]),
            }
        elif not validation["passed"]:
            result["status"] = "error"
            result["error"] = {"code": "validation_failed", "message": "Office validation failed"}
    elif operation == "create":
        if office_format == "docx":
            result["changes"] = _create_docx(payload, output_path)
        elif office_format == "xlsx":
            result["changes"] = _create_xlsx(payload, output_path)
        else:
            result["changes"] = _create_pptx(payload, output_path)
        result["artifact"] = _artifact(output_path)
        result["validation"] = _validate_openxml(output_path, office_format)
    else:
        assert input_path is not None
        if office_format == "docx":
            result["changes"] = _apply_docx(input_path, payload, output_path)
        elif office_format == "xlsx":
            result["changes"] = _apply_xlsx(input_path, payload, output_path)
        else:
            result["changes"] = _apply_pptx(input_path, payload, output_path)
        result["artifact"] = _artifact(output_path)
        result["validation"] = _validate_openxml(output_path, office_format)

    if input_path is not None and _sha256(input_path) != input_hash:
        raise OfficeRequestError("readonly_input_violated", "input artifact changed during operation")
    return result


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--request", required=True, type=Path)
    parser.add_argument("--result", required=True, type=Path)
    args = parser.parse_args()

    request: dict[str, Any] = {}
    started = time.monotonic()
    try:
        request = _read_json(args.request)
        result_path = args.result.expanduser().resolve()
        protected_paths = {
            path
            for value in (request.get("input_artifact"), request.get("output_artifact"))
            if isinstance(value, dict) and isinstance(value.get("path"), str)
            for path in (Path(value["path"]).expanduser().resolve(),)
        }
        if result_path in protected_paths or result_path == args.request.expanduser().resolve():
            print(
                json.dumps(
                    {
                        "status": "error",
                        "error": {
                            "code": "transport_path_collision",
                            "message": "result JSON path must differ from request and Office artifacts",
                        },
                    }
                ),
                file=sys.stderr,
            )
            raise SystemExit(2)
        result = execute(request)
        exit_code = 0 if result["status"] == "ok" else 2
    except OfficeRequestError as exc:
        result = _base_result(request.get("operation"), request.get("format"))
        result["status"] = exc.status
        result["error"] = {"code": exc.code, "message": str(exc)}
        exit_code = 2
    except Exception as exc:
        result = _base_result(request.get("operation"), request.get("format"))
        result["status"] = "error"
        result["error"] = {"code": "internal_error", "message": str(exc)}
        exit_code = 2
    result["duration_ms"] = round((time.monotonic() - started) * 1000, 3)
    _write_json_atomic(args.result, result)
    raise SystemExit(exit_code)


if __name__ == "__main__":
    main()
