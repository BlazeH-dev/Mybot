from __future__ import annotations

import hashlib
import json
import stat
import subprocess
import sys
import zipfile
from pathlib import Path
from typing import Any

import pytest
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from nanobot.agent.skills import SkillsLoader

ROOT = Path(__file__).resolve().parents[2]
SKILL_DIR = ROOT / "nanobot" / "skills" / "office-python"
SCRIPT = SKILL_DIR / "scripts" / "office.py"
FIXTURE_DIR = ROOT / "tests" / "fixtures" / "office_python"
SHARED_SCRIPTS = ROOT / "nanobot" / "skills" / "_shared" / "office_core" / "scripts"


def _sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _expand(value: Any, artifact_root: Path) -> Any:
    if isinstance(value, str):
        return value.replace("${ARTIFACT_ROOT}", str(artifact_root))
    if isinstance(value, list):
        return [_expand(item, artifact_root) for item in value]
    if isinstance(value, dict):
        return {key: _expand(item, artifact_root) for key, item in value.items()}
    return value


def _fixture_request(name: str, artifact_root: Path) -> dict[str, Any]:
    payload = json.loads((FIXTURE_DIR / name).read_text(encoding="utf-8"))
    return _expand(payload, artifact_root)


def _run_request(
    tmp_path: Path,
    request: dict[str, Any],
    *,
    name: str = "request",
    check: bool = True,
) -> tuple[subprocess.CompletedProcess[str], dict[str, Any]]:
    request_path = tmp_path / f"{name}.json"
    result_path = tmp_path / f"{name}.result.json"
    request_path.write_text(json.dumps(request), encoding="utf-8")
    process = subprocess.run(
        [sys.executable, str(SCRIPT), "--request", str(request_path), "--result", str(result_path)],
        cwd=ROOT,
        capture_output=True,
        text=True,
        check=False,
    )
    result = json.loads(result_path.read_text(encoding="utf-8"))
    if check and process.returncode != 0:
        raise AssertionError(f"office.py failed: {process.stderr}\n{result}")
    return process, result


def _create(tmp_path: Path, office_format: str) -> tuple[Path, dict[str, Any]]:
    artifact_root = tmp_path / "artifacts"
    request = _fixture_request(f"create_{office_format}.json", artifact_root)
    _, result = _run_request(tmp_path, request, name=f"create-{office_format}")
    return Path(result["artifact"]["path"]), result


def test_office_python_is_independently_discoverable() -> None:
    loader = SkillsLoader(ROOT)
    names = {entry["name"] for entry in loader.list_skills(filter_unavailable=False)}

    assert {"office-python", "officecli"}.issubset(names)
    assert "office-automation" not in names
    assert "_shared" not in names
    summary = loader.build_skills_summary()
    assert "**office-python**" in summary
    assert "**officecli**" in summary
    assert "General Python baseline" in summary
    assert str(SKILL_DIR / "SKILL.md") in summary

    disabled = SkillsLoader(ROOT, disabled_skills={"office-python"})
    configured = {entry["name"] for entry in disabled.list_skills(filter_unavailable=False)}
    assert "office-python" not in configured
    assert "officecli" in configured


def test_office_python_has_single_entrypoint_and_exact_constraints() -> None:
    manifest = (SKILL_DIR / "skill.yaml").read_text(encoding="utf-8")
    constraints = (SKILL_DIR / "references" / "constraints.txt").read_text(encoding="utf-8")

    assert "name: office-python" in manifest
    assert "scripts/office.py" in manifest
    assert "render_docx.py" not in manifest
    assert constraints.splitlines() == [
        "lxml==6.1.1",
        "openpyxl==3.1.5",
        "python-docx==1.2.0",
        "python-pptx==1.0.2",
    ]


@pytest.mark.parametrize("office_format", ["docx", "xlsx", "pptx"])
def test_create_inspect_query_apply_and_validate_are_neutral_and_readonly(
    tmp_path: Path,
    office_format: str,
) -> None:
    source, created = _create(tmp_path, office_format)
    artifact_root = source.parent
    source_hash = _sha256(source)

    assert created["status"] == "ok"
    assert created["validation"]["passed"] is True

    inspect_request = {
        "schema_version": 1,
        "operation": "inspect",
        "format": office_format,
        "input_artifact": {"path": str(source)},
        "selector": {"kind": "document"},
        "payload": {},
        "options": {"artifact_root": str(artifact_root)},
    }
    _, inspected = _run_request(tmp_path, inspect_request, name=f"inspect-{office_format}")
    assert inspected["matches"]

    if office_format == "docx":
        selector = {"kind": "paragraph", "text_contains": "Draft"}
        actions = [
            {"type": "set_text", "selector": selector, "value": "Final paragraph"},
            {"type": "set_header", "text": "Updated header"},
        ]
    elif office_format == "xlsx":
        selector = {"kind": "cell", "sheet": "Data", "range": "A1:B4"}
        actions = [
            {"type": "set_cell", "selector": {"sheet": "Data", "coordinate": "B2"}, "value": 15},
            {"type": "append_row", "selector": {"sheet": "Data"}, "values": ["Gamma", 30]},
        ]
    else:
        selector = {"kind": "shape", "slide": 1, "text_contains": "Draft slide"}
        actions = [
            {"type": "set_text", "selector": selector, "value": "Final slide"},
            {"type": "add_slide", "slide": {"layout": "title", "title": "Added slide"}},
        ]

    query_request = dict(inspect_request, operation="query", selector=selector)
    _, queried = _run_request(tmp_path, query_request, name=f"query-{office_format}")
    assert queried["matches"]
    if office_format == "xlsx":
        query_request["selector"] = {"kind": "cell", "sheet": "Data", "range": "B2"}
        _, single_cell = _run_request(tmp_path, query_request, name="query-xlsx-single-cell")
        assert single_cell["matches"][0]["coordinate"] == "B2"

    output = artifact_root / f"applied.{office_format}"
    apply_request = dict(
        inspect_request,
        operation="apply",
        output_artifact={"path": str(output)},
        payload={"actions": actions},
    )
    _, applied = _run_request(tmp_path, apply_request, name=f"apply-{office_format}")

    assert applied["status"] == "ok"
    assert applied["changes"]
    assert applied["validation"]["passed"] is True
    assert output.exists()
    assert _sha256(source) == source_hash

    validate_request = dict(
        inspect_request,
        operation="validate",
        input_artifact={"path": str(output)},
    )
    _, validated = _run_request(tmp_path, validate_request, name=f"validate-{office_format}")
    assert validated["validation"]["passed"] is True

    if office_format == "docx":
        document = Document(output)
        assert "Final paragraph" in [paragraph.text for paragraph in document.paragraphs]
        assert document.sections[0].header.paragraphs[0].text == "Updated header"
    elif office_format == "xlsx":
        workbook = load_workbook(output, data_only=False)
        assert workbook["Data"]["B2"].value == 15
        assert workbook["Data"]["B4"].value == "=SUM(B2:B3)"
        assert workbook["Data"]["A5"].value == "Gamma"
        assert len(workbook["Data"]._charts) == 1
        workbook.close()
    else:
        presentation = Presentation(output)
        assert len(presentation.slides) == 3
        assert presentation.slides[1].shapes.title.text == "Final slide"


def test_apply_is_atomic_when_a_later_action_fails(tmp_path: Path) -> None:
    source, _ = _create(tmp_path, "docx")
    output = source.parent / "existing.docx"
    output.write_bytes(b"existing-output")
    existing_hash = _sha256(output)
    request = {
        "schema_version": 1,
        "operation": "apply",
        "format": "docx",
        "input_artifact": {"path": str(source)},
        "output_artifact": {"path": str(output)},
        "selector": {"kind": "document"},
        "payload": {
            "actions": [
                {"type": "append_paragraph", "text": "would be temporary"},
                {"type": "set_text", "selector": {"kind": "paragraph", "index": 999}, "value": "no"},
            ]
        },
        "options": {"artifact_root": str(source.parent)},
    }

    process, result = _run_request(tmp_path, request, name="atomic-failure", check=False)

    assert process.returncode != 0
    assert result["error"]["code"] == "selector_no_match"
    assert _sha256(output) == existing_hash


def test_output_outside_artifact_root_is_rejected(tmp_path: Path) -> None:
    artifact_root = tmp_path / "artifacts"
    request = _fixture_request("create_docx.json", artifact_root)
    request["output_artifact"]["path"] = str(tmp_path / "outside.docx")

    process, result = _run_request(tmp_path, request, name="outside", check=False)

    assert process.returncode != 0
    assert result["error"]["code"] == "output_outside_artifact_root"
    assert not (tmp_path / "outside.docx").exists()


def test_result_json_cannot_overwrite_an_office_artifact(tmp_path: Path) -> None:
    source, _ = _create(tmp_path, "docx")
    source_hash = _sha256(source)
    request = {
        "schema_version": 1,
        "operation": "inspect",
        "format": "docx",
        "input_artifact": {"path": str(source)},
        "selector": {"kind": "document"},
        "payload": {},
        "options": {"artifact_root": str(source.parent)},
    }
    request_path = tmp_path / "collision-request.json"
    request_path.write_text(json.dumps(request), encoding="utf-8")

    process = subprocess.run(
        [sys.executable, str(SCRIPT), "--request", str(request_path), "--result", str(source)],
        cwd=ROOT,
        capture_output=True,
        text=True,
        check=False,
    )

    assert process.returncode != 0
    assert "transport_path_collision" in process.stderr
    assert _sha256(source) == source_hash


def test_tracked_changes_are_reported_as_unsupported(tmp_path: Path) -> None:
    source, _ = _create(tmp_path, "docx")
    rewritten = source.parent / "tracked.docx"
    with zipfile.ZipFile(source) as original, zipfile.ZipFile(rewritten, "w") as target:
        for item in original.infolist():
            content = original.read(item.filename)
            if item.filename == "word/document.xml":
                content = content.replace(b"<w:body>", b"<w:body><w:ins w:id=\"1\"><w:r><w:t>Tracked</w:t></w:r></w:ins>", 1)
            target.writestr(item, content)

    output = source.parent / "tracked-output.docx"
    request = {
        "schema_version": 1,
        "operation": "apply",
        "format": "docx",
        "input_artifact": {"path": str(rewritten)},
        "output_artifact": {"path": str(output)},
        "selector": {"kind": "document"},
        "payload": {"actions": [{"type": "append_paragraph", "text": "unsafe"}]},
        "options": {"artifact_root": str(source.parent)},
    }

    process, result = _run_request(tmp_path, request, name="tracked", check=False)

    assert process.returncode != 0
    assert result["status"] == "unsupported"
    assert result["error"]["code"] == "unsupported_features"
    assert "tracked_changes" in result["error"]["message"]
    assert not output.exists()


def test_render_requires_and_records_exact_external_libreoffice_version(tmp_path: Path) -> None:
    source, _ = _create(tmp_path, "pptx")
    fake_soffice = tmp_path / "soffice"
    fake_soffice.write_text(
        """#!/bin/sh
if [ "$1" = "--version" ]; then
  printf '%s\\n' 'LibreOffice 24.2.7.2'
  exit 0
fi
outdir=''
input=''
while [ "$#" -gt 0 ]; do
  if [ "$1" = "--outdir" ]; then
    shift
    outdir="$1"
  fi
  input="$1"
  shift
done
base=$(basename "$input")
base=${base%.*}
printf '%s' '%PDF-1.4 fixture' > "$outdir/$base.pdf"
""",
        encoding="utf-8",
    )
    fake_soffice.chmod(fake_soffice.stat().st_mode | stat.S_IXUSR)
    output = source.parent / "rendered.pdf"
    request = {
        "schema_version": 1,
        "operation": "render",
        "format": "pptx",
        "input_artifact": {"path": str(source)},
        "output_artifact": {"path": str(output)},
        "selector": {"kind": "presentation"},
        "payload": {},
        "options": {
            "artifact_root": str(source.parent),
            "libreoffice": {
                "path": str(fake_soffice),
                "expected_version": "LibreOffice 24.2.7.2",
            },
        },
    }

    _, result = _run_request(tmp_path, request, name="render")

    assert output.read_bytes().startswith(b"%PDF")
    assert result["validation"] == {
        "passed": True,
        "engine": "libreoffice",
        "path": str(fake_soffice),
        "version": "LibreOffice 24.2.7.2",
    }
    assert result["rendered_assets"][0]["sha256"] == _sha256(output)

    request["options"]["libreoffice"]["expected_version"] = "LibreOffice 25.0"
    process, mismatch = _run_request(tmp_path, request, name="render-mismatch", check=False)
    assert process.returncode != 0
    assert mismatch["error"]["code"] == "libreoffice_version_mismatch"


def test_shared_facts_and_inspector_work_with_neutral_xlsx_fixture(tmp_path: Path) -> None:
    source, _ = _create(tmp_path, "xlsx")
    schema_path = tmp_path / "schema.json"
    facts_path = tmp_path / "facts.json"
    spec_path = tmp_path / "metric.json"
    spec_path.write_text(
        json.dumps(
            {
                "schema_version": 1,
                "sheet": "Data",
                "metrics": [
                    {
                        "fact_id": "total_amount",
                        "name": "Total amount",
                        "calculation": "sum",
                        "column": "Amount",
                        "format": "integer",
                    }
                ],
            }
        ),
        encoding="utf-8",
    )

    subprocess.run(
        [sys.executable, str(SHARED_SCRIPTS / "inspect_workbook.py"), "--in", str(source), "--out", str(schema_path)],
        cwd=ROOT,
        check=True,
    )
    subprocess.run(
        [
            sys.executable,
            str(SHARED_SCRIPTS / "extract_facts.py"),
            "--in",
            str(source),
            "--spec",
            str(spec_path),
            "--out",
            str(facts_path),
        ],
        cwd=ROOT,
        check=True,
    )

    schema = json.loads(schema_path.read_text(encoding="utf-8"))
    fact = json.loads(facts_path.read_text(encoding="utf-8"))["facts"][0]
    assert schema["sheets"][0]["name"] == "Data"
    assert fact["fact_id"] == "total_amount"
    assert fact["value"] == 30
