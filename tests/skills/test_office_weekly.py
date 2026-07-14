from __future__ import annotations

import json
import os
import shutil
import subprocess
import sys
from pathlib import Path
from typing import Any

from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from nanobot.agent.skills import SkillsLoader

ROOT = Path(__file__).resolve().parents[2]
FIXTURE_DIR = ROOT / "tests" / "fixtures" / "office_weekly"
PYTHON_SKILL_DIR = ROOT / "nanobot" / "skills" / "office-automation"
OFFICECLI_SKILL_DIR = ROOT / "nanobot" / "skills" / "officecli"
SHARED_CORE_DIR = ROOT / "nanobot" / "skills" / "_shared" / "office_core"
PYTHON_SCRIPTS_DIR = PYTHON_SKILL_DIR / "scripts"
OFFICECLI_SCRIPTS_DIR = OFFICECLI_SKILL_DIR / "scripts"
SHARED_SCRIPTS_DIR = SHARED_CORE_DIR / "scripts"
PYTHON_REFERENCES_DIR = PYTHON_SKILL_DIR / "references"
OFFICECLI_REFERENCES_DIR = OFFICECLI_SKILL_DIR / "references"
SHARED_REFERENCES_DIR = SHARED_CORE_DIR / "references"


def _run_script(
    scripts_dir: Path,
    script_name: str,
    args: list[str | Path],
    *,
    check: bool = True,
) -> subprocess.CompletedProcess[str]:
    result = subprocess.run(
        [sys.executable, str(scripts_dir / script_name), *[str(arg) for arg in args]],
        cwd=ROOT,
        text=True,
        capture_output=True,
        check=False,
    )
    if check and result.returncode != 0:
        raise AssertionError(
            f"{script_name} failed with {result.returncode}\nSTDOUT:\n{result.stdout}\nSTDERR:\n{result.stderr}"
        )
    return result


def _read_json(path: Path) -> Any:
    return json.loads(path.read_text(encoding="utf-8"))


def _extract_fixture_facts(tmp_path: Path) -> Path:
    facts_path = tmp_path / "verified_facts.json"
    _run_script(
        SHARED_SCRIPTS_DIR,
        "extract_facts.py",
        [
            "--in",
            FIXTURE_DIR / "sales_data.xlsx",
            "--spec",
            SHARED_REFERENCES_DIR / "metric_spec.example.json",
            "--out",
            facts_path,
        ],
    )
    return facts_path


def test_independent_office_skills_are_discoverable_but_shared_core_is_not() -> None:
    loader = SkillsLoader(ROOT)

    entries = loader.list_skills(filter_unavailable=False)
    names = {entry["name"] for entry in entries}
    assert {"office-automation", "officecli"}.issubset(names)
    assert "_shared" not in names

    summary = loader.build_skills_summary()
    assert "**office-automation**" in summary
    assert "**officecli**" in summary
    assert "original Python" in summary
    assert "Default Office skill" in summary
    assert str(PYTHON_SKILL_DIR / "SKILL.md") in summary
    assert str(OFFICECLI_SKILL_DIR / "SKILL.md") in summary


def test_office_skill_availability_and_switches_are_independent(monkeypatch: Any) -> None:
    monkeypatch.setattr(
        "nanobot.agent.skills.shutil.which",
        lambda command: None if command == "officecli" else f"/usr/bin/{command}",
    )

    loader = SkillsLoader(ROOT)
    available_names = {entry["name"] for entry in loader.list_skills(filter_unavailable=True)}
    assert "office-automation" in available_names
    assert "officecli" not in available_names

    disabled_loader = SkillsLoader(ROOT, disabled_skills={"office-automation"})
    configured_names = {
        entry["name"] for entry in disabled_loader.list_skills(filter_unavailable=False)
    }
    assert "office-automation" not in configured_names
    assert "officecli" in configured_names


def test_officecli_runtime_contract_is_pinned() -> None:
    contract = _read_json(OFFICECLI_REFERENCES_DIR / "officecli-runtime.json")

    assert contract["provider"] == "officecli"
    assert contract["validated_version"] == "1.0.135"
    assert contract["allowed_batch_operations"] == ["add", "set"]
    assert {"raw-set", "plugins", "mcp", "watch", "install"}.issubset(
        contract["capabilities"]
    )
    assert "raw-set" in contract["policy_hints"]["ask"]
    assert all(len(asset["sha256"]) == 64 for asset in contract["assets"].values())


def test_inspect_workbook_emits_compact_schema(tmp_path: Path) -> None:
    output_path = tmp_path / "workbook_schema.json"

    _run_script(
        SHARED_SCRIPTS_DIR,
        "inspect_workbook.py",
        [
            "--in",
            FIXTURE_DIR / "sales_data.xlsx",
            "--out",
            output_path,
        ],
    )

    payload = _read_json(output_path)
    assert payload["schema_version"] == 1
    sheet = payload["sheets"][0]
    assert sheet["name"] == "sales_data"
    assert sheet["row_count"] == 8
    columns = {column["name"]: column for column in sheet["columns"]}
    assert columns["gmv_cny"]["type"] == "number"
    assert columns["region"]["type"] == "text"


def test_extract_facts_matches_expected_metrics(tmp_path: Path) -> None:
    facts_path = _extract_fixture_facts(tmp_path)
    payload = _read_json(facts_path)
    facts = {fact["fact_id"]: fact for fact in payload["facts"]}
    expected = {
        metric["id"]: metric
        for metric in _read_json(FIXTURE_DIR / "expected_metrics.json")["metrics"]
    }

    assert set(expected).issubset(facts)
    for fact_id, expected_metric in expected.items():
        fact = facts[fact_id]
        if isinstance(expected_metric["value"], float):
            assert abs(fact["value"] - expected_metric["value"]) < 1e-10
        else:
            assert fact["value"] == expected_metric["value"]
        assert fact["display_value"] == expected_metric["display_value"]
        assert fact["confidence"] == 1.0


def test_extract_facts_reports_missing_columns(tmp_path: Path) -> None:
    spec_path = tmp_path / "bad_metric_spec.json"
    spec_path.write_text(
        json.dumps(
            {
                "schema_version": 1,
                "sheet": "sales_data",
                "metrics": [
                    {
                        "fact_id": "missing",
                        "name": "Missing",
                        "calculation": "sum",
                        "column": "does_not_exist",
                        "format": "integer",
                    }
                ],
            }
        ),
        encoding="utf-8",
    )

    result = _run_script(
        SHARED_SCRIPTS_DIR,
        "extract_facts.py",
        [
            "--in",
            FIXTURE_DIR / "sales_data.xlsx",
            "--spec",
            spec_path,
            "--out",
            tmp_path / "facts.json",
        ],
        check=False,
    )

    assert result.returncode != 0
    assert "missing required column" in result.stderr


def test_extract_facts_treats_empty_numeric_cells_as_zero(tmp_path: Path) -> None:
    workbook_path = tmp_path / "with_empty.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "sales_data"
    sheet.append(["gmv_cny", "orders"])
    sheet.append([None, 1])
    sheet.append([100, 2])
    workbook.save(workbook_path)

    spec_path = tmp_path / "metric_spec.json"
    spec_path.write_text(
        json.dumps(
            {
                "schema_version": 1,
                "sheet": "sales_data",
                "metrics": [
                    {
                        "fact_id": "total_gmv_cny",
                        "name": "GMV total",
                        "calculation": "sum",
                        "column": "gmv_cny",
                        "unit": "CNY",
                        "format": "currency_cny",
                    }
                ],
            }
        ),
        encoding="utf-8",
    )
    facts_path = tmp_path / "facts.json"

    _run_script(
        SHARED_SCRIPTS_DIR,
        "extract_facts.py",
        ["--in", workbook_path, "--spec", spec_path, "--out", facts_path],
    )

    fact = _read_json(facts_path)["facts"][0]
    assert fact["value"] == 100
    assert fact["display_value"] == "CNY 100"


def test_validate_catches_unknown_fact_and_slide_limit(tmp_path: Path) -> None:
    facts_path = _extract_fixture_facts(tmp_path)
    invalid_slide_dsl = tmp_path / "invalid_slide_dsl.json"
    invalid_slide_dsl.write_text(
        json.dumps(
            {
                "schema_version": 1,
                "type": "slides",
                "slides": [
                    {
                        "title": f"Slide {index}",
                        "metrics": [{"label": "Bad", "fact_ref": "not_a_fact"}],
                    }
                    for index in range(7)
                ],
            }
        ),
        encoding="utf-8",
    )
    report_path = tmp_path / "quality_report.json"

    result = _run_script(
        PYTHON_SCRIPTS_DIR,
        "validate.py",
        [
            "--dsl",
            invalid_slide_dsl,
            "--facts",
            facts_path,
            "--constraints",
            FIXTURE_DIR / "expected_constraints.json",
            "--out",
            report_path,
        ],
        check=False,
    )

    assert result.returncode != 0
    report = _read_json(report_path)
    assert report["passed"] is False
    codes = {issue["code"] for issue in report["issues"]}
    assert {"too_many_slides", "unknown_fact_ref"}.issubset(codes)


def test_officecli_compiler_emits_bounded_replayable_commands(tmp_path: Path) -> None:
    facts_path = _extract_fixture_facts(tmp_path)
    report_batch = tmp_path / "report_batch.json"
    slide_batch = tmp_path / "slide_batch.json"

    _run_script(
        OFFICECLI_SCRIPTS_DIR,
        "compile_officecli.py",
        [
            "--kind",
            "docx",
            "--dsl",
            FIXTURE_DIR / "fixed_report_dsl.json",
            "--facts",
            facts_path,
            "--out",
            report_batch,
        ],
    )
    _run_script(
        OFFICECLI_SCRIPTS_DIR,
        "compile_officecli.py",
        [
            "--kind",
            "pptx",
            "--dsl",
            FIXTURE_DIR / "fixed_slide_dsl.json",
            "--facts",
            facts_path,
            "--out",
            slide_batch,
        ],
    )

    report_commands = _read_json(report_batch)
    slide_commands = _read_json(slide_batch)
    allowed_commands = {"add", "set"}
    assert report_commands
    assert slide_commands
    assert {command["command"] for command in report_commands} <= allowed_commands
    assert {command["command"] for command in slide_commands} <= allowed_commands
    assert all(command["command"] not in {"raw-set", "add-part"} for command in slide_commands)

    encoded = json.dumps([report_commands, slide_commands], ensure_ascii=False)
    assert "{{fact:" not in encoded
    assert "CNY 1,710,000" in encoded
    assert "7.81%" in encoded


def test_officecli_backend_real_binary(tmp_path: Path) -> None:
    officecli_bin = os.environ.get("OFFICECLI_TEST_BIN")
    if not officecli_bin:
        import pytest

        pytest.skip("set OFFICECLI_TEST_BIN to run the pinned OfficeCLI integration test")

    facts_path = _extract_fixture_facts(tmp_path)
    docx_path = tmp_path / "weekly_report.docx"
    pptx_path = tmp_path / "weekly_review.pptx"
    preview_dir = tmp_path / "previews"

    _run_script(
        OFFICECLI_SCRIPTS_DIR,
        "render_docx.py",
        [
            "--officecli-bin",
            officecli_bin,
            "--dsl",
            FIXTURE_DIR / "fixed_report_dsl.json",
            "--facts",
            facts_path,
            "--preview-dir",
            preview_dir,
            "--out",
            docx_path,
        ],
    )
    _run_script(
        OFFICECLI_SCRIPTS_DIR,
        "render_pptx.py",
        [
            "--officecli-bin",
            officecli_bin,
            "--dsl",
            FIXTURE_DIR / "fixed_slide_dsl.json",
            "--facts",
            facts_path,
            "--constraints",
            FIXTURE_DIR / "expected_constraints.json",
            "--preview-dir",
            preview_dir,
            "--out",
            pptx_path,
        ],
    )

    assert docx_path.exists()
    assert pptx_path.exists()
    assert docx_path.with_suffix(".docx.officecli-batch.json").exists()
    assert pptx_path.with_suffix(".pptx.officecli-batch.json").exists()
    assert docx_path.with_suffix(".docx.officecli-validation.json").exists()
    assert pptx_path.with_suffix(".pptx.officecli-validation.json").exists()
    assert docx_path.with_suffix(".docx.officecli-run.json").exists()
    assert pptx_path.with_suffix(".pptx.officecli-run.json").exists()
    assert list(preview_dir.glob("weekly_report*.png"))
    assert list(preview_dir.glob("weekly_review*.png"))

    docx_run = _read_json(docx_path.with_suffix(".docx.officecli-run.json"))
    pptx_run = _read_json(pptx_path.with_suffix(".pptx.officecli-run.json"))
    assert docx_run["engine_version"] == "1.0.135"
    assert pptx_run["engine_version"] == "1.0.135"
    assert docx_run["batch_sha256"]
    assert pptx_run["batch_sha256"]

    document = Document(docx_path)
    docx_text = "\n".join(
        [paragraph.text for paragraph in document.paragraphs]
        + [cell.text for table in document.tables for row in table.rows for cell in row.cells]
    )
    deck = Presentation(pptx_path)
    pptx_text = "\n".join(
        shape.text
        for slide in deck.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "CNY 1,710,000" in docx_text
    assert "CNY 1,710,000" in pptx_text
    assert "{{fact:" not in docx_text
    assert "{{fact:" not in pptx_text


def test_officecli_pptx_helper_enforces_slide_limit(tmp_path: Path) -> None:
    facts_path = _extract_fixture_facts(tmp_path)
    invalid_slide_dsl = tmp_path / "too_many_slides.json"
    invalid_slide_dsl.write_text(
        json.dumps(
            {
                "schema_version": 1,
                "type": "slides",
                "slides": [{"title": f"Slide {index}"} for index in range(7)],
            }
        ),
        encoding="utf-8",
    )

    result = _run_script(
        OFFICECLI_SCRIPTS_DIR,
        "render_pptx.py",
        [
            "--dsl",
            invalid_slide_dsl,
            "--facts",
            facts_path,
            "--constraints",
            FIXTURE_DIR / "expected_constraints.json",
            "--out",
            tmp_path / "too_many_slides.pptx",
        ],
        check=False,
    )

    assert result.returncode != 0
    assert "exceeds limit" in result.stderr


def test_office_weekly_deterministic_artifact_chain(tmp_path: Path) -> None:
    artifact_root = tmp_path / ".nanobot-runtime" / "artifacts" / "task_office_weekly"
    artifact_root.mkdir(parents=True)
    workbook_schema = artifact_root / "workbook_schema.json"
    facts_path = artifact_root / "verified_facts.json"
    report_dsl = artifact_root / "report_dsl.json"
    slide_dsl = artifact_root / "slide_dsl.json"
    plan_path = artifact_root / "plan.json"
    quality_report = artifact_root / "quality_report.json"
    docx_path = artifact_root / "weekly_report.docx"
    pptx_path = artifact_root / "weekly_review.pptx"

    shutil.copyfile(FIXTURE_DIR / "fixed_report_dsl.json", report_dsl)
    shutil.copyfile(FIXTURE_DIR / "fixed_slide_dsl.json", slide_dsl)
    shutil.copyfile(FIXTURE_DIR / "fixed_plan_done.json", plan_path)

    _run_script(
        SHARED_SCRIPTS_DIR,
        "inspect_workbook.py",
        ["--in", FIXTURE_DIR / "sales_data.xlsx", "--out", workbook_schema],
    )
    _run_script(
        SHARED_SCRIPTS_DIR,
        "extract_facts.py",
        [
            "--in",
            FIXTURE_DIR / "sales_data.xlsx",
            "--spec",
            SHARED_REFERENCES_DIR / "metric_spec.example.json",
            "--out",
            facts_path,
        ],
    )
    _run_script(
        PYTHON_SCRIPTS_DIR,
        "validate.py",
        [
            "--dsl",
            report_dsl,
            "--dsl",
            slide_dsl,
            "--facts",
            facts_path,
            "--constraints",
            FIXTURE_DIR / "expected_constraints.json",
            "--plan",
            plan_path,
            "--out",
            quality_report,
        ],
    )
    _run_script(
        PYTHON_SCRIPTS_DIR,
        "render_docx.py",
        [
            "--dsl",
            report_dsl,
            "--facts",
            facts_path,
            "--out",
            docx_path,
        ],
    )
    _run_script(
        PYTHON_SCRIPTS_DIR,
        "render_pptx.py",
        [
            "--dsl",
            slide_dsl,
            "--facts",
            facts_path,
            "--constraints",
            FIXTURE_DIR / "expected_constraints.json",
            "--out",
            pptx_path,
        ],
    )
    _run_script(
        PYTHON_SCRIPTS_DIR,
        "validate.py",
        [
            "--dsl",
            report_dsl,
            "--dsl",
            slide_dsl,
            "--facts",
            facts_path,
            "--constraints",
            FIXTURE_DIR / "expected_constraints.json",
            "--plan",
            plan_path,
            "--artifact-root",
            artifact_root,
            "--out",
            quality_report,
        ],
    )

    facts = {fact["fact_id"]: fact for fact in _read_json(facts_path)["facts"]}
    report = _read_json(quality_report)
    assert report["passed"] is True

    document = Document(docx_path)
    docx_text = "\n".join(
        [paragraph.text for paragraph in document.paragraphs]
        + [cell.text for table in document.tables for row in table.rows for cell in row.cells]
    )
    assert "{{fact:" not in docx_text
    for fact_id in ["total_gmv_cny", "total_orders", "conversion_rate", "top_region_by_gmv"]:
        assert facts[fact_id]["display_value"] in docx_text

    deck = Presentation(pptx_path)
    constraints = _read_json(FIXTURE_DIR / "expected_constraints.json")
    assert len(deck.slides) <= constraints["outputs"]["pptx_max_pages"]
    pptx_text = "\n".join(
        shape.text
        for slide in deck.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "{{fact:" not in pptx_text
    for fact_id in ["total_gmv_cny", "total_orders", "conversion_rate"]:
        assert facts[fact_id]["display_value"] in pptx_text
    for slide in deck.slides:
        assert any(shape.text.strip() for shape in slide.shapes if hasattr(shape, "text"))

    for artifact in constraints["outputs"]["required_artifacts"]:
        assert (artifact_root / artifact).exists()
